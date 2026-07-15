import asyncio
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
from typing import Any

from app.limits import DEFAULT_LIMITS, Limits


@dataclass(frozen=True)
class ExtractedDocument:
    text: str
    metadata: dict[str, Any] = field(default_factory=dict)


# Extensions/types nécessitant une extraction binaire (vs lecture texte directe).
_BINARY_EXTS = {".pdf", ".docx", ".xlsx", ".pptx"}
_BINARY_MIMES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def is_binary_document(filename: str, mime_type: str | None) -> bool:
    ext = Path(filename).suffix.lower()
    if ext in _BINARY_EXTS:
        return True
    return mime_type in _BINARY_MIMES


def _truncate(text: str, max_chars: int) -> tuple[str, bool, int]:
    """Tronque à `max_chars` caractères. Retourne (texte, truncated, total_chars)."""
    total = len(text)
    if total <= max_chars:
        return text, False, total
    return text[:max_chars], True, total


class DocumentExtractor(ABC):
    @abstractmethod
    async def extract(
        self,
        *,
        content: bytes,
        filename: str,
        mime_type: str | None,
    ) -> ExtractedDocument:
        """Extract plain text and metadata from a document."""


class LocalExtractor(DocumentExtractor):
    """Extraction légère hors-ligne pour documents digitaux.

    V1 : PDF (pdfplumber), Word .docx (python-docx), Excel .xlsx (openpyxl),
    PowerPoint .pptx (python-pptx). Les PDFs scannés (sans couche texte) et l'OCR
    sont hors V1.

    L'extraction applique des plafonds (pages, lignes, caractères) pour éviter
    d'envoyer un document gigantesque au LLM. Les métadonnées indiquent la
    troncature (`truncated`, `chars_total`, `pages_total`, etc.).
    """

    def __init__(self, limits: Limits = DEFAULT_LIMITS) -> None:
        self._limits = limits

    async def extract(
        self,
        *,
        content: bytes,
        filename: str,
        mime_type: str | None,
    ) -> ExtractedDocument:
        ext = Path(filename).suffix.lower()
        if ext == ".pdf" or mime_type == "application/pdf":
            return await asyncio.to_thread(self._extract_pdf, content, filename)
        if ext == ".docx":
            return await asyncio.to_thread(self._extract_docx, content, filename)
        if ext == ".xlsx":
            return await asyncio.to_thread(self._extract_xlsx, content, filename)
        if ext == ".pptx":
            return await asyncio.to_thread(self._extract_pptx, content, filename)
        return ExtractedDocument(text="", metadata={"note": "unsupported", "filename": filename})

    def _finalize(
        self,
        text: str,
        *,
        filename: str,
        extractor: str,
        extra: dict[str, Any],
    ) -> ExtractedDocument:
        max_chars = self._limits.extract_max_chars
        truncated_text, truncated, chars_total = _truncate(text, max_chars)
        metadata: dict[str, Any] = {
            "filename": filename,
            "extractor": extractor,
            "chars_total": chars_total,
            "chars_returned": len(truncated_text),
            "truncated": truncated,
            **extra,
        }
        return ExtractedDocument(text=truncated_text, metadata=metadata)

    def _extract_pdf(self, content: bytes, filename: str) -> ExtractedDocument:
        import pdfplumber  # lazy import (lourd)

        max_pages = self._limits.extract_max_pages
        pages_text: list[str] = []
        total_pages = 0
        with pdfplumber.open(BytesIO(content)) as pdf:
            total_pages = len(pdf.pages)
            for page in pdf.pages[:max_pages]:
                pages_text.append(page.extract_text() or "")
        text = "\n\n".join(pages_text).strip()
        pages_read = min(total_pages, max_pages)
        return self._finalize(
            text,
            filename=filename,
            extractor="pdfplumber",
            extra={
                "pages_total": total_pages,
                "pages_read": pages_read,
                "pages_truncated": total_pages > max_pages,
            },
        )

    def _extract_docx(self, content: bytes, filename: str) -> ExtractedDocument:
        from docx import Document  # python-docx

        document = Document(BytesIO(content))
        # On capte par budget caractères plutôt que par nombre de paragraphes
        # (les paragraphes n'ont pas de notion de "page").
        max_chars = self._limits.extract_max_chars
        kept: list[str] = []
        total_paragraphs = 0
        used = 0
        for p in document.paragraphs:
            if not (p.text and p.text.strip()):
                continue
            total_paragraphs += 1
            piece = p.text.strip()
            if used + len(piece) > max_chars:
                remaining = max_chars - used
                if remaining > 0:
                    kept.append(piece[:remaining])
                break
            kept.append(piece)
            used += len(piece)
        text = "\n\n".join(kept).strip()
        truncated_paragraphs = total_paragraphs > len(kept)
        return self._finalize(
            text,
            filename=filename,
            extractor="python-docx",
            extra={
                "paragraphs_total": total_paragraphs,
                "paragraphs_read": len(kept),
                "paragraphs_truncated": truncated_paragraphs,
            },
        )

    def _extract_xlsx(self, content: bytes, filename: str) -> ExtractedDocument:
        from openpyxl import load_workbook

        workbook = load_workbook(BytesIO(content), read_only=True, data_only=True)
        max_rows = self._limits.extract_max_rows
        lines: list[str] = []
        sheets_total = len(workbook.sheetnames)
        rows_truncated = False
        for sheet in workbook.worksheets:
            lines.append(f"# {sheet.title}")
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                if not any(cell is not None for cell in row):
                    continue
                if row_count >= max_rows:
                    rows_truncated = True
                    break
                lines.append("\t".join("" if cell is None else str(cell) for cell in row))
                row_count += 1
        workbook.close()
        text = "\n".join(lines).strip()
        return self._finalize(
            text,
            filename=filename,
            extractor="openpyxl",
            extra={
                "sheets_total": sheets_total,
                "max_rows_per_sheet": max_rows,
                "rows_truncated": rows_truncated,
            },
        )

    def _extract_pptx(self, content: bytes, filename: str) -> ExtractedDocument:
        from pptx import Presentation

        presentation = Presentation(BytesIO(content))
        max_slides = self._limits.extract_max_pages
        slides_text: list[str] = []
        total_slides = len(presentation.slides)
        for index, slide in enumerate(presentation.slides, start=1):
            if index > max_slides:
                break
            parts: list[str] = [f"# Slide {index}"]
            for shape in slide.shapes:
                text = getattr(shape, "text", None)
                if text and text.strip():
                    parts.append(text.strip())
            slides_text.append("\n".join(parts))
        text = "\n\n".join(slides_text).strip()
        slides_read = min(total_slides, max_slides)
        return self._finalize(
            text,
            filename=filename,
            extractor="python-pptx",
            extra={
                "slides_total": total_slides,
                "slides_read": slides_read,
                "slides_truncated": total_slides > max_slides,
            },
        )
