"""Aperçu HTML simplifié des documents Office et PDF."""

from __future__ import annotations

import html
from io import BytesIO
from pathlib import Path
from typing import Any

from app.documents.extractor import LocalExtractor
from app.limits import DEFAULT_LIMITS, Limits

_PREVIEW_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".gif"}
_PREVIEW_TEXT_EXTS = {
    ".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm", ".log",
}


def preview_type_for_path(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".docx":
        return "docx"
    if ext == ".xlsx":
        return "xlsx"
    if ext == ".pptx":
        return "pptx"
    if ext == ".pdf":
        return "pdf"
    if ext in _PREVIEW_IMAGE_EXTS:
        return "image"
    if ext in _PREVIEW_TEXT_EXTS:
        return "text"
    return "unsupported"


async def render_preview(
    path: Path,
    *,
    limits: Limits = DEFAULT_LIMITS,
) -> dict[str, Any]:
    """Produit un aperçu HTML léger pour le panneau droit."""
    preview_type = preview_type_for_path(path)
    title = path.name

    if preview_type == "unsupported":
        return {"type": "unsupported", "title": title, "html": ""}

    if preview_type == "image":
        return {"type": "image", "title": title, "html": ""}

    content = path.read_bytes()
    if len(content) > limits.extract_max_input_bytes:
        return {
            "type": preview_type,
            "title": title,
            "html": (
                f"<p>{html.escape('File too large for preview')}</p>"
            ),
        }

    if preview_type == "docx":
        return {
            "type": "docx",
            "title": title,
            "html": _render_docx_html(content),
        }

    if preview_type == "xlsx":
        return {
            "type": "xlsx",
            "title": title,
            "html": _render_xlsx_html(content),
        }

    if preview_type == "pptx":
        try:
            pptx_html = _render_pptx_html(content, limits=limits)
        except Exception:  # noqa: BLE001 - fichier corrompu / non-pptx
            pptx_html = f"<p>{html.escape('Preview unavailable for this PowerPoint file')}</p>"
        return {
            "type": "pptx",
            "title": title,
            "html": pptx_html,
        }

    if preview_type == "pdf":
        return {
            "type": "pdf",
            "title": title,
            "html": _render_pdf_html(content, limits=limits),
        }

    text = content.decode("utf-8", errors="replace")
    if len(text) > limits.extract_max_chars:
        text = text[: limits.extract_max_chars]
    return {
        "type": "text",
        "title": title,
        "html": f"<pre>{html.escape(text)}</pre>",
    }


def _render_pdf_html(content: bytes, *, limits: Limits = DEFAULT_LIMITS) -> str:
    extractor = LocalExtractor(limits=limits)
    extracted = extractor._extract_pdf(content, "preview.pdf")
    text = html.escape(extracted.text or "")
    return f"<pre>{text}</pre>"


def _render_docx_html(content: bytes) -> str:
    from docx import Document
    from docx.text.paragraph import Paragraph

    document = Document(BytesIO(content))
    parts: list[str] = []

    def paragraph_html(paragraph: Paragraph) -> str:
        text = html.escape(paragraph.text.strip())
        if not text:
            return ""
        style_name = (paragraph.style.name or "").lower() if paragraph.style else ""
        if "heading 1" in style_name:
            return f"<h1>{text}</h1>"
        if "heading 2" in style_name:
            return f"<h2>{text}</h2>"
        if paragraph.style and paragraph.style.name == "List Bullet":
            return f"<li>{text}</li>"
        return f"<p>{text}</p>"

    in_list = False
    for paragraph in document.paragraphs:
        piece = paragraph_html(paragraph)
        if not piece:
            continue
        if piece.startswith("<li>"):
            if not in_list:
                parts.append("<ul>")
                in_list = True
            parts.append(piece)
        else:
            if in_list:
                parts.append("</ul>")
                in_list = False
            parts.append(piece)
    if in_list:
        parts.append("</ul>")

    return "\n".join(parts) if parts else "<p></p>"


def _render_xlsx_html(content: bytes) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(BytesIO(content), read_only=True, data_only=True)
    sheet = workbook.worksheets[0]
    max_rows = 50
    max_cols = 20

    rows_html: list[str] = []
    for row_index, row in enumerate(sheet.iter_rows(values_only=True)):
        if row_index >= max_rows:
            break
        cells = row[:max_cols]
        if not any(cell is not None for cell in cells):
            continue
        cell_html = "".join(
            f"<td>{html.escape('' if cell is None else str(cell))}</td>"
            for cell in cells
        )
        rows_html.append(f"<tr>{cell_html}</tr>")

    workbook.close()
    if not rows_html:
        return "<table></table>"
    return "<table>" + "".join(rows_html) + "</table>"


def _render_pptx_html(content: bytes, *, limits: Limits = DEFAULT_LIMITS) -> str:
    from pptx import Presentation

    presentation = Presentation(BytesIO(content))
    max_slides = limits.extract_max_pages
    parts: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        if index > max_slides:
            break
        slide_parts: list[str] = [
            f'<div class="wp-pptx-slide"><h2>Slide {index}</h2>',
        ]
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if not text or not text.strip():
                continue
            for line in text.strip().splitlines():
                stripped = line.strip()
                if stripped:
                    slide_parts.append(f"<p>{html.escape(stripped)}</p>")
        slide_parts.append("</div>")
        parts.extend(slide_parts)

    return "\n".join(parts) if parts else "<p></p>"
