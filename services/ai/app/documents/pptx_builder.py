"""Génération native PowerPoint (.pptx) avec thèmes et layouts structurés."""

from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from typing import Any

# Layouts supportés par write_pptx (catalogue produit, pas d'improvisation libre).
PPTX_LAYOUTS = frozenset(
    {
        "title",
        "section",
        "bullets",
        "two_column",
        "kpi_row",
        "quote",
        "closing",
    }
)
PPTX_THEMES = frozenset({"light", "dark", "improba"})
MAX_PPTX_SLIDES = 60
MAX_KPI_CARDS = 4

_SLIDE_W = 13.333  # inches, 16:9 widescreen
_SLIDE_H = 7.5


@dataclass(frozen=True)
class _Theme:
    name: str
    bg: tuple[int, int, int]
    surface: tuple[int, int, int]
    text: tuple[int, int, int]
    muted: tuple[int, int, int]
    accent: tuple[int, int, int]
    on_accent: tuple[int, int, int]


_THEMES: dict[str, _Theme] = {
    "light": _Theme(
        name="light",
        bg=(255, 255, 255),
        surface=(247, 244, 239),
        text=(28, 25, 23),
        muted=(87, 83, 78),
        accent=(196, 163, 90),
        on_accent=(28, 25, 23),
    ),
    "dark": _Theme(
        name="dark",
        bg=(28, 25, 23),
        surface=(41, 37, 36),
        text=(245, 245, 244),
        muted=(168, 162, 158),
        accent=(212, 175, 55),
        on_accent=(28, 25, 23),
    ),
    "improba": _Theme(
        name="improba",
        bg=(250, 248, 244),
        surface=(255, 255, 255),
        text=(32, 61, 82),
        muted=(90, 110, 125),
        accent=(196, 163, 90),
        on_accent=(32, 61, 82),
    ),
}


def build_pptx_bytes(
    *,
    slides: list[dict[str, Any]] | None = None,
    theme: str = "improba",
    fidelity: str = "editable",
) -> bytes:
    """Construit un fichier PowerPoint 16:9 à partir de slides structurées."""
    fidelity_key = (fidelity or "editable").strip().lower()
    if fidelity_key == "visual":
        from app.documents.pptx_svg import build_pptx_visual_bytes

        content, _meta = build_pptx_visual_bytes(slides, theme=theme)
        return content
    content, _meta = build_pptx_editable_bytes(slides=slides, theme=theme)
    return content


def build_pptx_editable_bytes(
    *,
    slides: list[dict[str, Any]] | None = None,
    theme: str = "improba",
    skip_critique: bool = False,
) -> tuple[bytes, dict[str, Any]]:
    """Rendu natif éditable (python-pptx) après normalisation + critique."""
    from pptx import Presentation
    from pptx.util import Inches

    from app.documents.slides_critique import critique_and_fix

    theme_key = (theme or "improba").strip().lower()
    if theme_key not in _THEMES:
        allowed = ", ".join(sorted(PPTX_THEMES))
        raise ValueError(f"Thème inconnu {theme!r}. Utilise: {allowed}")
    palette = _THEMES[theme_key]

    if skip_critique:
        normalized = list(slides or [])
        if len(normalized) > MAX_PPTX_SLIDES:
            raise ValueError(
                f"Trop de slides ({len(normalized)}). Maximum: {MAX_PPTX_SLIDES}"
            )
        critique_issues = 0
    else:
        critique = critique_and_fix(slides, theme=theme_key)
        normalized = critique.slides
        critique_issues = len(critique.issues)

    presentation = Presentation()
    presentation.slide_width = Inches(_SLIDE_W)
    presentation.slide_height = Inches(_SLIDE_H)

    for index, slide_data in enumerate(normalized):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        _fill_background(slide, palette.bg)
        _draw_grammar(slide, slide_data=slide_data, theme=palette, index=index)

    buf = BytesIO()
    presentation.save(buf)
    return buf.getvalue(), {"critique_issues": critique_issues}


def _rgb(color: tuple[int, int, int]):
    from pptx.dml.color import RGBColor

    return RGBColor(*color)


def _fill_background(slide: Any, color: tuple[int, int, int]) -> None:
    """Fond plein. Ajouté en premier shape = arrière-plan (pas de réordonnancement XML)."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(_SLIDE_W),
        Inches(_SLIDE_H),
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)


def _add_rect(
    slide: Any,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: tuple[int, int, int],
) -> Any:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    return shape


def _add_textbox(
    slide: Any,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int,
    color: tuple[int, int, int],
    bold: bool = False,
    align: str = "left",
) -> None:
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    run.font.name = "Calibri"
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT


def _add_bullets(
    slide: Any,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    items: list[str],
    color: tuple[int, int, int],
    font_size: int = 22,
) -> None:
    from pptx.util import Inches, Pt

    cleaned = [item.strip() for item in items if item and item.strip()]
    if not cleaned:
        return

    box = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    tf = box.text_frame
    tf.word_wrap = True
    for index, text in enumerate(cleaned):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = _rgb(color)
        p.font.name = "Calibri"
        p.space_after = Pt(10)


def _as_str_list(value: Any) -> list[str]:
    if not isinstance(value, list):
        return []
    out: list[str] = []
    for item in value:
        if isinstance(item, str) and item.strip():
            out.append(item.strip())
        elif isinstance(item, dict):
            label = item.get("label") or item.get("title") or ""
            val = item.get("value") or item.get("text") or ""
            if label or val:
                out.append(f"{label}: {val}".strip(": ").strip())
    return out


def _hierarchy_fields(slide_data: dict[str, Any]) -> tuple[str, list[str], list[str]]:
    hierarchy = (
        slide_data.get("hierarchy") if isinstance(slide_data.get("hierarchy"), dict) else {}
    )
    primary = str(hierarchy.get("primary") or "").strip()
    secondary = hierarchy.get("secondary") if isinstance(hierarchy.get("secondary"), list) else []
    secondary_s = [str(s).strip() for s in secondary if str(s).strip()]
    tertiary = hierarchy.get("tertiary") if isinstance(hierarchy.get("tertiary"), list) else []
    tertiary_s = [str(s).strip() for s in tertiary if str(s).strip()]
    return primary, secondary_s, tertiary_s


def _visual_fields(slide_data: dict[str, Any]) -> tuple[str, list[Any]]:
    visual = slide_data.get("visual") if isinstance(slide_data.get("visual"), dict) else {}
    vtype = str(visual.get("type") or "none")
    items = visual.get("items") if isinstance(visual.get("items"), list) else []
    return vtype, items


def _split_visual_items(
    items: list[Any], secondary: list[str]
) -> tuple[list[str], list[str], str, str]:
    left_t = right_t = ""
    left: list[str] = []
    right: list[str] = []
    for item in items:
        if not isinstance(item, dict):
            continue
        side = str(item.get("side") or "").lower()
        title = str(item.get("title") or "").strip()
        lines = item.get("lines") if isinstance(item.get("lines"), list) else []
        lines_s = [str(x).strip() for x in lines if str(x).strip()]
        if side == "right":
            right_t = title or right_t
            right.extend(lines_s)
        else:
            left_t = title or left_t
            left.extend(lines_s)
    if not left and not right and secondary:
        mid = max(1, len(secondary) // 2)
        left, right = secondary[:mid], secondary[mid:]
    return left, right, left_t, right_t


def _kpi_from_visual(items: list[Any], secondary: list[str]) -> list[dict[str, str]]:
    cards: list[dict[str, str]] = []
    for item in items:
        if isinstance(item, dict):
            cards.append(
                {
                    "label": str(item.get("label") or "").strip(),
                    "value": str(item.get("value") or "-").strip(),
                }
            )
        if len(cards) >= MAX_KPI_CARDS:
            break
    if not cards:
        for line in secondary[:MAX_KPI_CARDS]:
            if ":" in line:
                label, _, value = line.partition(":")
                cards.append({"label": label.strip(), "value": value.strip()})
            else:
                cards.append({"label": "", "value": line})
    return cards


def _draw_grammar(
    slide: Any,
    *,
    slide_data: dict[str, Any],
    theme: _Theme,
    index: int,
) -> None:
    grammar = str(slide_data.get("grammar") or "sequence")
    primary, secondary, tertiary = _hierarchy_fields(slide_data)
    vtype, items = _visual_fields(slide_data)
    bullets = secondary + tertiary

    _add_rect(slide, left=0, top=0, width=0.18, height=_SLIDE_H, fill=theme.accent)

    if grammar == "hero":
        title = primary or "Présentation"
        combined = bullets
        subtitle = combined[0] if combined else ""
        extra_bullets = combined[1:] if len(combined) > 1 else []
        _draw_title_slide(slide, title=title, subtitle=subtitle, theme=theme)
        if extra_bullets:
            _add_bullets(
                slide,
                left=0.9,
                top=5.2,
                width=11.2,
                height=1.8,
                items=extra_bullets,
                color=theme.muted,
                font_size=18,
            )
        return

    if grammar in {"split", "comparison"}:
        left_items, right_items, left_title, right_title = _split_visual_items(
            items, bullets
        )
        _draw_two_column_slide(
            slide,
            title=primary,
            left_title=left_title,
            right_title=right_title,
            left_items=left_items,
            right_items=right_items,
            theme=theme,
        )
        return

    if grammar == "dashboard" or vtype == "kpi_row":
        cards = _kpi_from_visual(items, bullets)
        _draw_kpi_slide(slide, title=primary, cards=cards, theme=theme)
        return

    if grammar == "editorial" or vtype == "quote":
        quote = primary
        attribution = secondary[0] if secondary else ""
        _draw_quote_slide(
            slide, quote=quote, attribution=attribution, theme=theme
        )
        return

    if grammar == "sequence":
        _draw_bullets_slide(
            slide, title=primary or f"Slide {index + 1}", bullets=bullets, theme=theme
        )
        return

    if grammar == "diagram":
        steps = bullets or _lines_from_visual_items(items)
        _draw_bullets_slide(
            slide, title=primary or f"Slide {index + 1}", bullets=steps, theme=theme
        )
        return

    _draw_bullets_slide(
        slide, title=primary or f"Slide {index + 1}", bullets=bullets, theme=theme
    )


def _lines_from_visual_items(items: list[Any]) -> list[str]:
    out: list[str] = []
    for item in items:
        if isinstance(item, str) and item.strip():
            out.append(item.strip())
        elif isinstance(item, dict):
            text = item.get("text") or item.get("label") or item.get("value")
            if text:
                out.append(str(text).strip())
    return out


def _draw_title_slide(
    slide: Any,
    *,
    title: str,
    subtitle: str,
    theme: _Theme,
) -> None:
    _add_rect(
        slide,
        left=0.18,
        top=0,
        width=_SLIDE_W - 0.18,
        height=0.12,
        fill=theme.accent,
    )
    _add_textbox(
        slide,
        left=0.9,
        top=2.4,
        width=11.2,
        height=1.6,
        text=title,
        font_size=44,
        color=theme.text,
        bold=True,
    )
    if subtitle:
        _add_textbox(
            slide,
            left=0.9,
            top=4.2,
            width=11.2,
            height=1.0,
            text=subtitle,
            font_size=22,
            color=theme.muted,
        )


def _draw_section_slide(slide: Any, *, title: str, subtitle: str, theme: _Theme) -> None:
    _add_textbox(
        slide,
        left=0.9,
        top=2.6,
        width=11.2,
        height=1.2,
        text=title,
        font_size=40,
        color=theme.text,
        bold=True,
        align="center",
    )
    _add_rect(slide, left=5.4, top=4.0, width=2.5, height=0.08, fill=theme.accent)
    if subtitle:
        _add_textbox(
            slide,
            left=1.5,
            top=4.4,
            width=10.3,
            height=1.0,
            text=subtitle,
            font_size=20,
            color=theme.muted,
            align="center",
        )


def _draw_quote_slide(
    slide: Any, *, quote: str, attribution: str, theme: _Theme
) -> None:
    body = quote or "Citation"
    _add_textbox(
        slide,
        left=1.4,
        top=2.2,
        width=10.4,
        height=2.4,
        text=f"« {body} »",
        font_size=28,
        color=theme.text,
        align="center",
    )
    if attribution:
        _add_textbox(
            slide,
            left=1.4,
            top=4.8,
            width=10.4,
            height=0.6,
            text=f"- {attribution}",
            font_size=18,
            color=theme.muted,
            align="center",
        )


def _draw_kpi_slide(
    slide: Any,
    *,
    title: str,
    cards: list[dict[str, str]],
    theme: _Theme,
) -> None:
    if title:
        _add_textbox(
            slide,
            left=0.9,
            top=0.55,
            width=11.2,
            height=0.8,
            text=title,
            font_size=32,
            color=theme.text,
            bold=True,
        )
        _add_rect(slide, left=0.9, top=1.35, width=1.8, height=0.07, fill=theme.accent)
    count = max(len(cards), 1)
    gap = 0.35
    usable = 11.4
    card_w = (usable - gap * (count - 1)) / count
    top = 2.4
    for i, card in enumerate(cards):
        left = 0.9 + i * (card_w + gap)
        _add_rect(slide, left=left, top=top, width=card_w, height=2.8, fill=theme.surface)
        _add_rect(slide, left=left, top=top, width=card_w, height=0.1, fill=theme.accent)
        value = str(card.get("value") or "-").strip()
        label = str(card.get("label") or "").strip()
        _add_textbox(
            slide,
            left=left + 0.2,
            top=top + 0.7,
            width=card_w - 0.4,
            height=1.0,
            text=value,
            font_size=28 if len(value) < 12 else 22,
            color=theme.text,
            bold=True,
            align="center",
        )
        if label:
            _add_textbox(
                slide,
                left=left + 0.2,
                top=top + 1.8,
                width=card_w - 0.4,
                height=0.6,
                text=label,
                font_size=14,
                color=theme.muted,
                align="center",
            )


def _draw_two_column_slide(
    slide: Any,
    *,
    title: str,
    left_title: str,
    right_title: str,
    left_items: list[str],
    right_items: list[str],
    theme: _Theme,
) -> None:
    if title:
        _add_textbox(
            slide,
            left=0.9,
            top=0.45,
            width=11.2,
            height=0.7,
            text=title,
            font_size=30,
            color=theme.text,
            bold=True,
        )
        _add_rect(slide, left=0.9, top=1.15, width=1.8, height=0.07, fill=theme.accent)
    col_top = 1.5
    col_top_left = col_top + (0.55 if left_title else 0)
    col_top_right = col_top + (0.55 if right_title else 0)
    if left_title:
        _add_textbox(
            slide,
            left=0.9,
            top=col_top,
            width=5.3,
            height=0.5,
            text=left_title,
            font_size=18,
            color=theme.accent,
            bold=True,
        )
    if right_title:
        _add_textbox(
            slide,
            left=6.9,
            top=col_top,
            width=5.3,
            height=0.5,
            text=right_title,
            font_size=18,
            color=theme.accent,
            bold=True,
        )
    _add_bullets(
        slide,
        left=0.9,
        top=col_top_left,
        width=5.3,
        height=4.8,
        items=left_items,
        color=theme.text,
        font_size=18,
    )
    _add_bullets(
        slide,
        left=6.9,
        top=col_top_right,
        width=5.3,
        height=4.8,
        items=right_items,
        color=theme.text,
        font_size=18,
    )


def _draw_bullets_slide(
    slide: Any,
    *,
    title: str,
    bullets: list[str],
    theme: _Theme,
    subtitle: str = "",
) -> None:
    if title:
        _add_textbox(
            slide,
            left=0.9,
            top=0.45,
            width=11.2,
            height=0.7,
            text=title,
            font_size=30,
            color=theme.text,
            bold=True,
        )
        _add_rect(slide, left=0.9, top=1.15, width=1.8, height=0.07, fill=theme.accent)
    bullets_top = 1.5
    if subtitle:
        _add_textbox(
            slide,
            left=0.9,
            top=1.5,
            width=11.2,
            height=0.6,
            text=subtitle,
            font_size=18,
            color=theme.muted,
        )
        bullets_top = 2.3
    if bullets:
        _add_bullets(
            slide,
            left=0.9,
            top=bullets_top,
            width=11.2,
            height=5.2,
            items=bullets,
            color=theme.text,
            font_size=20,
        )


def _draw_layout(
    slide: Any,
    *,
    layout: str,
    data: dict[str, Any],
    theme: _Theme,
    index: int,
) -> None:
    """Compat tests internes : délègue au rendu par grammaire."""
    from app.documents.slides_schema import normalize_slide

    slide_data = normalize_slide({**data, "layout": layout}, index=index + 1)
    _draw_grammar(slide, slide_data=slide_data, theme=theme, index=index)
