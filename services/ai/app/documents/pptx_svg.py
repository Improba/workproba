"""PPTX visuel (screenshot HTML) avec repli sur rendu natif éditable."""

from __future__ import annotations

from io import BytesIO
from typing import Any

from app.documents.slides_chromium import chromium_available, html_slides_to_pngs
from app.documents.slides_critique import critique_and_fix
from app.documents.slides_html import render_deck_html

_SLIDE_W = 13.333
_SLIDE_H = 7.5


def build_pptx_visual_bytes(
    slides: list[dict[str, Any]] | None,
    theme: str = "improba",
) -> tuple[bytes, dict[str, Any]]:
    """Construit un PPTX image (une slide PNG par page) ou repli éditable.

    Retourne (bytes, metadata) avec ``render_mode`` = ``visual`` ou ``editable_fallback``.
    Sans Chromium disponible ou en cas d'échec capture, repli sur le builder natif.
    """
    from app.documents.pptx_builder import build_pptx_editable_bytes

    critique = critique_and_fix(slides, theme=theme)
    normalized = critique.slides
    meta: dict[str, Any] = {
        "fidelity": "visual",
        "theme": theme,
        "critique_issues": len(critique.issues),
    }

    if not chromium_available():
        meta["render_mode"] = "editable_fallback"
        meta["fallback_reason"] = "chromium_unavailable"
        content, _ = build_pptx_editable_bytes(
            slides=normalized, theme=theme, skip_critique=True
        )
        return content, meta

    html = render_deck_html(normalized, theme=theme)
    try:
        pngs = html_slides_to_pngs(html)
    except Exception as exc:  # noqa: BLE001 - repli non bloquant
        meta["render_mode"] = "editable_fallback"
        meta["fallback_reason"] = f"chromium_error:{type(exc).__name__}"
        content, _ = build_pptx_editable_bytes(
            slides=normalized, theme=theme, skip_critique=True
        )
        return content, meta

    if not pngs:
        meta["render_mode"] = "editable_fallback"
        meta["fallback_reason"] = "no_slides_captured"
        content, _ = build_pptx_editable_bytes(
            slides=normalized, theme=theme, skip_critique=True
        )
        return content, meta

    meta["render_mode"] = "visual"
    return _build_pptx_from_pngs(pngs), meta


def _build_pptx_from_pngs(pngs: list[bytes]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    presentation.slide_width = Inches(_SLIDE_W)
    presentation.slide_height = Inches(_SLIDE_H)
    blank_layout = presentation.slide_layouts[6]

    for png in pngs:
        slide = presentation.slides.add_slide(blank_layout)
        slide.shapes.add_picture(BytesIO(png), Inches(0), Inches(0), width=Inches(_SLIDE_W))

    buf = BytesIO()
    presentation.save(buf)
    return buf.getvalue()


__all__ = ["build_pptx_visual_bytes"]
