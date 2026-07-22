"""Tests outils Office natifs (write_docx/xlsx/pptx/pdf)."""

from __future__ import annotations

from base64 import b64decode
from io import BytesIO
from pathlib import Path

import pytest
from pptx import Presentation
from pydantic_ai.models.test import TestModel

from app.agent.tools import ToolDeps, build_agent
from app.documents.writer import (
    build_docx_bytes,
    build_pdf_bytes,
    build_pptx_bytes,
    build_xlsx_bytes,
    require_path_extension,
)
from app.limits import DEFAULT_LIMITS
from app.sandbox.runner import SandboxRunner

from conftest import FakeProjectClient


def test_build_docx_bytes_contains_paragraphs() -> None:
    content = build_docx_bytes(title="Titre", paragraphs=["Ligne 1", "Ligne 2"])
    assert content.startswith(b"PK")
    assert len(content) > 100


def test_build_xlsx_bytes_contains_sheet() -> None:
    content = build_xlsx_bytes(
        sheets=[{"name": "Data", "rows": [["A", "B"], [1, 2]]}]
    )
    assert content.startswith(b"PK")


def test_build_pdf_bytes_produces_pdf() -> None:
    content = build_pdf_bytes(
        title="Rapport",
        sections=[{"heading": "Intro", "body": "Contenu"}],
    )
    assert content.startswith(b"%PDF")


def _pptx_all_text(content: bytes) -> str:
    presentation = Presentation(BytesIO(content))
    texts = [
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "text", None)
    ]
    return "\n".join(texts)


def test_build_pptx_bytes_split_includes_tertiary() -> None:
    content = build_pptx_bytes(
        slides=[
            {
                "grammar": "comparison",
                "hierarchy": {
                    "primary": "Choix",
                    "secondary": ["Option A", "Option B"],
                    "tertiary": ["Detail A", "Detail B"],
                },
            }
        ]
    )
    joined = _pptx_all_text(content)
    assert "Option A" in joined
    assert "Detail A" in joined
    assert "Detail B" in joined


def test_build_pptx_bytes_hero_includes_tertiary() -> None:
    content = build_pptx_bytes(
        slides=[
            {
                "grammar": "hero",
                "hierarchy": {
                    "primary": "Titre hero",
                    "secondary": ["Sous-titre"],
                    "tertiary": ["Puce extra"],
                },
            }
        ]
    )
    joined = _pptx_all_text(content)
    assert "Titre hero" in joined
    assert "Sous-titre" in joined
    assert "Puce extra" in joined


def test_build_pdf_fallback_from_slides_includes_tertiary(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    from app.documents.writer import _build_pdf_fallback_from_slides

    captured: list[str] = []
    from reportlab.platypus import Paragraph as RealParagraph

    def _capture(text: str, style: object) -> RealParagraph:
        captured.append(text)
        return RealParagraph(text, style)

    monkeypatch.setattr("reportlab.platypus.Paragraph", _capture)

    content = _build_pdf_fallback_from_slides(
        [
            {
                "hierarchy": {
                    "primary": "Slide PDF",
                    "secondary": ["Ligne sec"],
                    "tertiary": ["Ligne ter"],
                },
            }
        ]
    )
    assert content.startswith(b"%PDF")
    joined = "\n".join(captured)
    assert "Slide PDF" in joined
    assert "Ligne sec" in joined
    assert "Ligne ter" in joined


def test_build_pdf_bytes_from_slides_uses_fallback_without_chromium(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    from app.documents.writer import build_pdf_bytes_from_slides

    monkeypatch.setattr(
        "app.documents.slides_chromium.chromium_available", lambda: False
    )
    content = build_pdf_bytes_from_slides(
        slides=[
            {
                "grammar": "sequence",
                "hierarchy": {
                    "primary": "Slide PDF",
                    "secondary": ["Ligne sec"],
                    "tertiary": ["Ligne ter"],
                },
            }
        ]
    )
    assert content.startswith(b"%PDF")
    assert len(content) > 100


def test_build_pptx_visual_fallback_render_mode(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    from app.documents.pptx_svg import build_pptx_visual_bytes

    monkeypatch.setattr("app.documents.pptx_svg.chromium_available", lambda: False)
    content, meta = build_pptx_visual_bytes(
        slides=[{"grammar": "hero", "hierarchy": {"primary": "Visuel", "secondary": []}}]
    )
    assert content.startswith(b"PK")
    assert meta["render_mode"] == "editable_fallback"
    assert meta["fallback_reason"] == "chromium_unavailable"


@pytest.mark.asyncio
async def test_write_pdf_slides_smoke(tmp_path: Path) -> None:
    from pydantic_ai import RunContext

    client = FakeProjectClient()
    deps = _deps(tmp_path, client)
    agent = build_agent(TestModel())
    tool = agent._function_toolset.tools["write_pdf"]
    ctx = RunContext(deps=deps, model=TestModel(), usage=None, prompt=None, tool_call_id="tc5")

    result = await tool.function(
        ctx,
        path="deck/slides.pdf",
        slides=[
            {
                "grammar": "sequence",
                "hierarchy": {
                    "primary": "Deck",
                    "secondary": ["Point"],
                    "tertiary": ["Detail"],
                },
            }
        ],
    )
    assert result.get("name") == "deck/slides.pdf"
    raw = b64decode(client.saved[0]["content_base64"])
    assert raw.startswith(b"%PDF")


def test_build_pptx_bytes_grammar_hero() -> None:
    content = build_pptx_bytes(
        slides=[
            {
                "grammar": "hero",
                "hierarchy": {"primary": "Semantic", "secondary": ["Tag"]},
            }
        ]
    )
    assert content.startswith(b"PK")
    presentation = Presentation(BytesIO(content))
    assert len(presentation.slides) == 1
    texts = [
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "text", None)
    ]
    assert any("Semantic" in text for text in texts)


def test_build_pptx_bytes_visual_fidelity_fallback_without_chromium() -> None:
    content = build_pptx_bytes(
        fidelity="visual",
        slides=[{"layout": "title", "title": "Visual fallback", "subtitle": "ok"}],
    )
    assert content.startswith(b"PK")
    presentation = Presentation(BytesIO(content))
    assert len(presentation.slides) >= 1


def test_build_pptx_bytes_is_real_presentation() -> None:
    content = build_pptx_bytes(
        theme="improba",
        slides=[
            {"layout": "title", "title": "Lapin", "subtitle": "Pitch"},
            {
                "layout": "bullets",
                "title": "Plan",
                "bullets": ["A", "B", "C"],
            },
            {
                "layout": "kpi_row",
                "title": "Chiffres",
                "metrics": [
                    {"label": "ARR", "value": "1.2M"},
                    {"label": "NPS", "value": "72"},
                ],
            },
            {
                "layout": "two_column",
                "title": "Comparaison",
                "left_title": "Avant",
                "right_title": "Après",
                "left": ["Lent"],
                "right": ["Rapide"],
            },
            {
                "layout": "quote",
                "quote": "On livre du vrai pptx",
                "attribution": "Improba",
            },
            {"layout": "section", "title": "Suite"},
            {
                "layout": "closing",
                "title": "Merci",
                "subtitle": "Questions ?",
                "bullets": ["Contact"],
            },
        ],
    )
    assert content.startswith(b"PK")
    presentation = Presentation(BytesIO(content))
    assert len(presentation.slides) == 7
    texts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text:
                texts.append(text)
    joined = "\n".join(texts)
    assert "Lapin" in joined
    assert "1.2M" in joined
    assert "On livre du vrai pptx" in joined


def test_build_pptx_bytes_empty_slides_uses_default_title() -> None:
    for slides in (None, []):
        content = build_pptx_bytes(slides=slides)
        presentation = Presentation(BytesIO(content))
        assert len(presentation.slides) == 1
        texts = [
            shape.text
            for slide in presentation.slides
            for shape in slide.shapes
            if getattr(shape, "text", None)
        ]
        assert any("Présentation" in text for text in texts)


def test_build_pptx_bytes_rejects_non_dict_slides() -> None:
    with pytest.raises(ValueError, match="dict"):
        build_pptx_bytes(slides=["pas un dict"])  # type: ignore[arg-type]


def test_build_pptx_bytes_rejects_invalid_theme() -> None:
    with pytest.raises(ValueError, match="Thème"):
        build_pptx_bytes(
            theme="neon",
            slides=[{"layout": "title", "title": "X"}],
        )


def test_build_pptx_bytes_rejects_invalid_layout() -> None:
    with pytest.raises(ValueError, match="layout"):
        build_pptx_bytes(slides=[{"layout": "carousel", "title": "X"}])


def test_build_pptx_bytes_rejects_empty_bullets_slide() -> None:
    with pytest.raises(ValueError, match="bullets"):
        build_pptx_bytes(slides=[{"layout": "bullets"}])


def test_build_pptx_bytes_rejects_empty_two_column() -> None:
    with pytest.raises(ValueError, match="two_column"):
        build_pptx_bytes(slides=[{"layout": "two_column", "title": "T"}])


def test_build_pptx_bytes_rejects_too_many_kpis() -> None:
    with pytest.raises(ValueError, match="maximum 4"):
        build_pptx_bytes(
            slides=[
                {
                    "layout": "kpi_row",
                    "title": "KPIs",
                    "metrics": [
                        {"label": f"K{i}", "value": str(i)} for i in range(5)
                    ],
                }
            ]
        )


def test_build_pptx_bytes_dashboard_includes_tertiary_without_items() -> None:
    content = build_pptx_bytes(
        slides=[
            {
                "grammar": "dashboard",
                "hierarchy": {
                    "primary": "Tableau",
                    "secondary": ["Revenu: 10"],
                    "tertiary": ["Marge: 20"],
                },
                "visual": {"type": "kpi_row", "items": []},
            }
        ]
    )
    joined = _pptx_all_text(content)
    assert "Revenu" in joined
    assert "10" in joined
    assert "Marge" in joined
    assert "20" in joined


def test_build_pptx_editable_skip_critique_rejects_too_many_slides() -> None:
    from app.documents.pptx_builder import MAX_PPTX_SLIDES, build_pptx_editable_bytes

    slides = [
        {"grammar": "hero", "hierarchy": {"primary": f"S{i}"}}
        for i in range(MAX_PPTX_SLIDES + 1)
    ]
    with pytest.raises(ValueError, match="Trop de slides"):
        build_pptx_editable_bytes(slides=slides, skip_critique=True)


def test_build_pptx_bytes_rejects_too_many_slides() -> None:
    from app.documents.pptx_builder import MAX_PPTX_SLIDES

    slides = [
        {"layout": "title", "title": f"S{i}"} for i in range(MAX_PPTX_SLIDES + 1)
    ]
    with pytest.raises(ValueError, match="Trop de slides"):
        build_pptx_bytes(slides=slides)


def test_require_path_extension_rejects_mismatch() -> None:
    with pytest.raises(ValueError, match=r"\.pptx"):
        require_path_extension("deck.docx", ".pptx")
    require_path_extension("deck.pptx", ".pptx")


def test_write_tools_registered() -> None:
    agent = build_agent(TestModel())
    names = sorted(agent._function_toolset.tools.keys())
    assert "write_docx" in names
    assert "write_xlsx" in names
    assert "write_pptx" in names
    assert "write_pdf" in names


def _deps(tmp_path: Path, client: FakeProjectClient) -> ToolDeps:
    return ToolDeps(
        context=__import__("app.agent.tools", fromlist=["ToolContext"]).ToolContext(
            tenant_id="t",
            project_id="p",
            session_id="s",
            documents=[],
            project_root=tmp_path,
        ),
        project_client=client,
        sandbox_runner=SandboxRunner(timeout_seconds=30, limits=DEFAULT_LIMITS),
        limits=DEFAULT_LIMITS,
        confirmation_gate=None,
    )


@pytest.mark.asyncio
async def test_write_docx_persists_file(tmp_path: Path) -> None:
    from pydantic_ai import RunContext

    client = FakeProjectClient()
    deps = _deps(tmp_path, client)
    agent = build_agent(TestModel())
    tool = agent._function_toolset.tools["write_docx"]

    ctx = RunContext(deps=deps, model=TestModel(), usage=None, prompt=None, tool_call_id="tc1")
    result = await tool.function(ctx, path="out/report.docx", title="R", paragraphs=["Hi"])
    assert result.get("name") == "out/report.docx"
    assert client.saved
    assert client.saved[0]["mime_type"].endswith("document")


@pytest.mark.asyncio
async def test_write_pptx_editable_reports_critique_issues(tmp_path: Path) -> None:
    from pydantic_ai import RunContext

    client = FakeProjectClient()
    deps = _deps(tmp_path, client)
    agent = build_agent(TestModel())
    tool = agent._function_toolset.tools["write_pptx"]
    ctx = RunContext(deps=deps, model=TestModel(), usage=None, prompt=None, tool_call_id="tc2c")

    secondary = [f"Point {i} avec du texte supplementaire" for i in range(7)]
    await tool.function(
        ctx,
        path="deck/dense.pptx",
        slides=[
            {
                "grammar": "sequence",
                "hierarchy": {"primary": "Dense", "secondary": secondary},
            }
        ],
    )
    saved_meta = client.saved[0]["metadata"]
    assert saved_meta["critique_issues"] > 0
    assert saved_meta["fidelity"] == "editable"


@pytest.mark.asyncio
async def test_write_pptx_persists_real_pptx(tmp_path: Path) -> None:
    from pydantic_ai import RunContext
    from pydantic_ai.exceptions import ModelRetry

    client = FakeProjectClient()
    deps = _deps(tmp_path, client)
    agent = build_agent(TestModel())
    tool = agent._function_toolset.tools["write_pptx"]
    ctx = RunContext(deps=deps, model=TestModel(), usage=None, prompt=None, tool_call_id="tc2")

    result = await tool.function(
        ctx,
        path="deck/pitch.pptx",
        theme="dark",
        slides=[{"layout": "title", "title": "Hello", "subtitle": "World"}],
    )
    assert result.get("name") == "deck/pitch.pptx"
    assert client.saved[0]["mime_type"].endswith("presentation")
    raw = b64decode(client.saved[0]["content_base64"])
    assert Presentation(BytesIO(raw)).slides

    with pytest.raises(ModelRetry):
        await tool.function(
            ctx,
            path="deck/fake.docx",
            slides=[{"layout": "title", "title": "Nope"}],
        )


@pytest.mark.asyncio
async def test_write_pptx_rejects_invalid_theme(tmp_path: Path) -> None:
    from pydantic_ai import RunContext
    from pydantic_ai.exceptions import ModelRetry

    client = FakeProjectClient()
    deps = _deps(tmp_path, client)
    agent = build_agent(TestModel())
    tool = agent._function_toolset.tools["write_pptx"]
    ctx = RunContext(deps=deps, model=TestModel(), usage=None, prompt=None, tool_call_id="tc2b")

    with pytest.raises(ModelRetry, match="Thème"):
        await tool.function(
            ctx,
            path="deck.pptx",
            theme="neon",
            slides=[{"layout": "title", "title": "X"}],
        )
    assert not client.saved


@pytest.mark.asyncio
async def test_write_docx_rejects_pptx_extension(tmp_path: Path) -> None:
    from pydantic_ai import RunContext
    from pydantic_ai.exceptions import ModelRetry

    client = FakeProjectClient()
    deps = _deps(tmp_path, client)
    agent = build_agent(TestModel())
    tool = agent._function_toolset.tools["write_docx"]
    ctx = RunContext(deps=deps, model=TestModel(), usage=None, prompt=None, tool_call_id="tc3")

    with pytest.raises(ModelRetry, match=r"\.docx"):
        await tool.function(ctx, path="evil.pptx", title="X", paragraphs=["y"])
    assert not client.saved


@pytest.mark.asyncio
async def test_generate_document_rejects_office_extension(tmp_path: Path) -> None:
    from pydantic_ai import RunContext
    from pydantic_ai.exceptions import ModelRetry

    client = FakeProjectClient()
    deps = _deps(tmp_path, client)
    agent = build_agent(TestModel())
    tool = agent._function_toolset.tools["generate_document"]
    ctx = RunContext(deps=deps, model=TestModel(), usage=None, prompt=None, tool_call_id="tc4")

    with pytest.raises(ModelRetry, match="write_pptx"):
        await tool.function(
            ctx,
            name="deck.pptx",
            mime_type="text/markdown",
            content_markdown="# Fake",
        )
    assert not client.saved
