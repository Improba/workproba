"""Tests endpoint /documents/preview."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

import pytest
from docx import Document
from fastapi.testclient import TestClient
from openpyxl import Workbook

import app.auth as authmod
import app.main as mainmod
from app.documents.preview import _render_pptx_html


@pytest.fixture(autouse=True)
def _loopback(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setattr(authmod, "is_loopback_host", lambda host: True)


def _headers() -> dict[str, str]:
    return {"X-Internal-Secret": "desktop-dev-secret"}


def test_preview_docx(tmp_path: Path) -> None:
    doc = Document()
    doc.add_heading("Titre", level=1)
    doc.add_paragraph("Corps")
    path = tmp_path / "note.docx"
    doc.save(path)

    with TestClient(mainmod.app) as client:
        resp = client.get(
            "/documents/preview",
            params={
                "path": "note.docx",
                "workspace_data_dir": str(tmp_path),
                "project_path": str(tmp_path),
            },
            headers=_headers(),
        )
    assert resp.status_code == 200
    data = resp.json()
    assert data["type"] == "docx"
    assert "<h1>Titre</h1>" in data["html"]
    assert "Corps" in data["html"]


def test_preview_xlsx(tmp_path: Path) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["A", "B"])
    ws.append([1, 2])
    path = tmp_path / "data.xlsx"
    wb.save(path)

    with TestClient(mainmod.app) as client:
        resp = client.get(
            "/documents/preview",
            params={
                "path": "data.xlsx",
                "workspace_data_dir": str(tmp_path),
                "project_path": str(tmp_path),
            },
            headers=_headers(),
        )
    assert resp.status_code == 200
    data = resp.json()
    assert data["type"] == "xlsx"
    assert "<table>" in data["html"]
    assert "A" in data["html"]


def test_preview_rejects_prefix_traversal(tmp_path: Path) -> None:
    """Bloque /tmp/foo-bar quand base=/tmp/foo (startswith insuffisant)."""
    base = tmp_path / "foo"
    base.mkdir()
    sibling = tmp_path / "foo-bar"
    sibling.mkdir()
    secret = sibling / "secret.txt"
    secret.write_text("secret", encoding="utf-8")

    with TestClient(mainmod.app) as client:
        resp = client.get(
            "/documents/preview",
            params={
                "path": "../foo-bar/secret.txt",
                "workspace_data_dir": str(tmp_path),
                "project_path": str(base),
            },
            headers=_headers(),
        )
    assert resp.status_code == 403


def test_preview_rejects_path_traversal(tmp_path: Path) -> None:
    outside = tmp_path.parent / "outside.txt"
    outside.write_text("secret", encoding="utf-8")

    with TestClient(mainmod.app) as client:
        resp = client.get(
            "/documents/preview",
            params={
                "path": "../outside.txt",
                "workspace_data_dir": str(tmp_path),
                "project_path": str(tmp_path),
            },
            headers=_headers(),
        )
    assert resp.status_code == 403


def test_preview_pptx(tmp_path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Titre"
    slide.placeholders[1].text = "Contenu slide"
    path = tmp_path / "deck.pptx"
    prs.save(path)

    with TestClient(mainmod.app) as client:
        resp = client.get(
            "/documents/preview",
            params={
                "path": "deck.pptx",
                "workspace_data_dir": str(tmp_path),
                "project_path": str(tmp_path),
            },
            headers=_headers(),
        )
    assert resp.status_code == 200
    data = resp.json()
    assert data["type"] == "pptx"
    assert "Titre" in data["html"]
    assert "Contenu slide" in data["html"]


def test_preview_pptx_empty_deck(tmp_path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    path = tmp_path / "empty.pptx"
    prs.save(path)

    with TestClient(mainmod.app) as client:
        resp = client.get(
            "/documents/preview",
            params={
                "path": "empty.pptx",
                "workspace_data_dir": str(tmp_path),
                "project_path": str(tmp_path),
            },
            headers=_headers(),
        )
    assert resp.status_code == 200
    data = resp.json()
    assert data["type"] == "pptx"
    assert "Slide 1" in data["html"]
    assert "<p>" not in data["html"]


def test_preview_pptx_truncates_slides(tmp_path: Path) -> None:
    from pptx import Presentation

    from app.limits import Limits

    prs = Presentation()
    title_layout = prs.slide_layouts[1]
    for index in range(55):
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = f"Slide {index + 1}"
    path = tmp_path / "long.pptx"
    prs.save(path)

    with TestClient(mainmod.app) as client:
        resp = client.get(
            "/documents/preview",
            params={
                "path": "long.pptx",
                "workspace_data_dir": str(tmp_path),
                "project_path": str(tmp_path),
            },
            headers=_headers(),
        )
    assert resp.status_code == 200
    html = resp.json()["html"]
    assert "Slide 50" in html
    assert "Slide 51" not in html

    truncated = _render_pptx_html(
        path.read_bytes(),
        limits=Limits(extract_max_pages=3),
    )
    assert "Slide 3" in truncated
    assert "Slide 4" not in truncated


def test_preview_separate_project_and_data_dirs(tmp_path: Path) -> None:
    """Desktop : folder_path et data_dir (.workproba) sont des chemins distincts."""
    data_dir = tmp_path / "data"
    project_dir = tmp_path / "project"
    data_dir.mkdir()
    project_dir.mkdir()
    doc = Document()
    doc.add_paragraph("Hello")
    doc.save(project_dir / "note.docx")

    with TestClient(mainmod.app) as client:
        resp = client.get(
            "/documents/preview",
            params={
                "path": "note.docx",
                "workspace_data_dir": str(data_dir),
                "project_path": str(project_dir),
            },
            headers=_headers(),
        )
    assert resp.status_code == 200
    assert "Hello" in resp.json()["html"]


def test_preview_image_type(tmp_path: Path) -> None:
    img = tmp_path / "pic.png"
    img.write_bytes(b"\x89PNG\r\n\x1a\n")

    with TestClient(mainmod.app) as client:
        resp = client.get(
            "/documents/preview",
            params={
                "path": "pic.png",
                "workspace_data_dir": str(tmp_path),
                "project_path": str(tmp_path),
            },
            headers=_headers(),
        )
    assert resp.status_code == 200
    assert resp.json()["type"] == "image"
