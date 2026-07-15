"""Tests extraction documents (hors réseau)."""

from __future__ import annotations

import asyncio
from io import BytesIO

from app.documents.extractor import LocalExtractor, is_binary_document


def test_is_binary_document_by_extension() -> None:
    assert is_binary_document("rapport.pdf", None)
    assert is_binary_document("note.docx", None)
    assert is_binary_document("data.xlsx", None)
    assert is_binary_document("slides.pptx", None)


def test_is_binary_document_by_mime() -> None:
    assert is_binary_document("file.bin", "application/pdf")
    assert not is_binary_document("file.bin", "text/plain")
    assert not is_binary_document("readme.md", None)


def test_extract_docx() -> None:
    from docx import Document

    doc = Document()
    doc.add_paragraph("Bonjour Workproba")
    doc.add_paragraph("Seconde ligne")
    buf = BytesIO()
    doc.save(buf)

    extractor = LocalExtractor()
    result = asyncio.run(
        extractor.extract(content=buf.getvalue(), filename="note.docx", mime_type=None)
    )
    assert "Bonjour Workproba" in result.text
    assert "Seconde ligne" in result.text


def test_extract_pptx() -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Intro"
    slide.placeholders[1].text = "Points clés"
    buf = BytesIO()
    prs.save(buf)

    extractor = LocalExtractor()
    result = asyncio.run(
        extractor.extract(content=buf.getvalue(), filename="deck.pptx", mime_type=None)
    )
    assert "Intro" in result.text
    assert "Points clés" in result.text
    assert result.metadata.get("slides_total") == 1


def test_extract_unsupported_falls_back_gracefully() -> None:
    extractor = LocalExtractor()
    result = asyncio.run(
        extractor.extract(content=b"plain text", filename="note.txt", mime_type="text/plain")
    )
    # .txt n'est pas un format binaire géré ; l'extracteur renvoie un texte vide
    # (la lecture texte directe est gérée côté LocalProjectClient).
    assert result.text == ""
